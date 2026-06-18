from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

GREEN = RGBColor(0x22, 0xC5, 0x5E)
DARK_GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x2D, 0x1F)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xF5)
GRAY = RGBColor(0x6B, 0x7B, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_green_bar(slide, left=Inches(0.6), top=Inches(0.8), height=Inches(0.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

def add_title_text(slide, text, left=Inches(0.85), top=Inches(0.7), width=Inches(10), size=Pt(36)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN

def add_body_text(slide, text, left=Inches(0.85), top=Inches(1.8), width=Inches(5.5), size=Pt(18), color=DARK, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    return tf

def add_card(slide, left, top, width, height, title, desc, icon_text="", border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(1)

    if icon_text:
        icon_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(0.6), Inches(0.6))
        itf = icon_box.text_frame
        ip = itf.paragraphs[0]
        ip.text = icon_text
        ip.font.size = Pt(24)

    title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.9), width - Inches(0.6), Inches(0.5))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = DARK_GREEN

    desc_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.4), width - Inches(0.6), height - Inches(1.8))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(13)
    dp.font.color.rgb = GRAY

# ========== SLIDE 0: COUVERTURE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_GREEN

logo_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(4), Inches(0.8))
tf = logo_box.text_frame
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "Eco"
run1.font.size = Pt(32)
run1.font.bold = True
run1.font.color.rgb = WHITE
run2 = p.add_run()
run2.text = "Track"
run2.font.size = Pt(32)
run2.font.bold = True
run2.font.color.rgb = ORANGE

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7), Inches(2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Gestion Durable des"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.text = "Dechets Urbains"
p2.font.size = Pt(44)
p2.font.bold = True
p2.font.color.rgb = WHITE

sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(7), Inches(1))
tf = sub_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Plateforme connectee et optimisee pour transformer l'assainissement urbain au Cameroun."
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(6), Inches(0.6))
tf = badge_box.text_frame
p = tf.paragraphs[0]
p.text = "Orange Summer Challenge 2026 - Douala"
p.font.size = Pt(16)
p.font.color.rgb = YELLOW

# Pilote box
pilote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2.5), Inches(3.5), Inches(3))
pilote.fill.solid()
pilote.fill.fore_color.rgb = RGBColor(0x1A, 0x4A, 0x2A)
pilote.line.fill.background()
pb = slide.shapes.add_textbox(Inches(9.3), Inches(3.2), Inches(3), Inches(2))
tf = pb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Projet Pilote"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Specifications Techniques et Fonctionnelles"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p2.alignment = PP_ALIGN.CENTER

# ========== SLIDE 1: LE PROBLEME ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "LE PROBLEME : LA CRISE DES DECHETS")

stat_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(5.5), Inches(1.2))
stat_shape.fill.solid()
stat_shape.fill.fore_color.rgb = DARK_GREEN
stat_shape.line.fill.background()

stat_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(2), Inches(1))
tf = stat_box.text_frame
p = tf.paragraphs[0]
p.text = "70%"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = ORANGE

stat_desc = slide.shapes.add_textbox(Inches(3.2), Inches(2.2), Inches(3), Inches(1))
tf = stat_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "des dechets solides restes non collectes au Cameroun."
p.font.size = Pt(16)
p.font.color.rgb = WHITE

# Consequences
cons_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.8), Inches(5.5), Inches(1.5))
tf = cons_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Consequences directes :"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
p2 = tf.add_paragraph()
p2.text = "Les caniveaux bouches provoquent de graves inondations recurrentes a Douala et Yaounde, empechant de nombreux enfants d'acceder a l'ecole."
p2.font.size = Pt(14)
p2.font.color.rgb = DARK

# Logistique
log_box = slide.shapes.add_textbox(Inches(0.85), Inches(5.5), Inches(5.5), Inches(1.2))
tf = log_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Logistique aveugle : "
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = DARK
run2 = p.add_run()
run2.text = "Les camions d'Hysacam effectuent des tournees sans aucune donnee en temps reel, entrainant la proliferation de decharges sauvages."
run2.font.size = Pt(14)
run2.font.color.rgb = DARK

# ========== SLIDE 2: LA SOLUTION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "LA SOLUTION : L'ECOSYSTEME ECOTRACK")

desc_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(11), Inches(0.8))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Une plateforme "
run.font.size = Pt(18)
run.font.color.rgb = DARK
run2 = p.add_run()
run2.text = "tripartite innovante"
run2.font.size = Pt(18)
run2.font.color.rgb = GREEN
run2.font.bold = True
run3 = p.add_run()
run3.text = " connectant les acteurs cles de l'environnement urbain pour une efficacite absolue :"
run3.font.size = Pt(18)
run3.font.color.rgb = DARK

add_card(slide, Inches(0.85), Inches(3.0), Inches(3.5), Inches(3.2),
         "Citoyens Actifs",
         "Les citoyens se transforment en capteurs urbains en signalant et geolocalisant les depots via leur smartphone en quelques secondes.",
         icon_text="\U0001F46B")

add_card(slide, Inches(4.85), Inches(3.0), Inches(3.5), Inches(3.2),
         "Collecteurs Optimises",
         "Les agents recoivent des feuilles de route dynamiques calculees pour minimiser les trajets d'intervention.",
         icon_text="\U0001F69A")

add_card(slide, Inches(8.85), Inches(3.0), Inches(3.5), Inches(3.2),
         "Communes Avisees",
         "Les administrations centralisent les alertes, pilotent les equipes de collecte et visualisent la reduction de l'empreinte carbone.",
         icon_text="\U0001F3DB")

# ========== SLIDE 3: SIGNALEMENT CITOYEN ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "SIGNALEMENT CITOYEN")

sub_title = slide.shapes.add_textbox(Inches(0.85), Inches(1.6), Inches(5), Inches(0.5))
tf = sub_title.text_frame
p = tf.paragraphs[0]
p.text = "En 30 secondes chrono"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ORANGE

desc = slide.shapes.add_textbox(Inches(0.85), Inches(2.3), Inches(5.5), Inches(2))
tf = desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Grace a un systeme de geolocalisation simple, l'utilisateur prend une photo et signale un point noir. La position GPS exacte et l'image sont envoyees instantanement sur notre base de donnees."
p.font.size = Pt(16)
p.font.color.rgb = DARK

quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(4.5), Inches(5.5), Inches(1.2))
quote_shape.fill.solid()
quote_shape.fill.fore_color.rgb = RGBColor(0xEC, 0xFC, 0xEC)
quote_shape.line.fill.background()
quote_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.7), Inches(5), Inches(1))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Chaque signalement permet d\'actualiser la carte dynamique consultee par les communes et partenaires logistiques."'
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = DARK_GREEN

# Phone mockup on right
phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.5), Inches(3.5), Inches(5))
phone.fill.solid()
phone.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xF0)
phone.line.color.rgb = GREEN
phone.line.width = Pt(3)
phone_text = slide.shapes.add_textbox(Inches(8.5), Inches(3.5), Inches(2.5), Inches(1.5))
tf = phone_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "\U0001F4F1 Signaler"
p.font.size = Pt(24)
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\U0001F4CD GPS + Photo"
p2.font.size = Pt(16)
p2.font.color.rgb = GRAY
p2.alignment = PP_ALIGN.CENTER

# ========== SLIDE 4: ITINERAIRES OPTIMISES ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "ITINERAIRES LOGISTIQUES OPTIMISES")

algo_title = slide.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6), Inches(0.6))
tf = algo_title.text_frame
p = tf.paragraphs[0]
p.text = "Algorithme TSP & Carbone"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ORANGE

algo_desc = slide.shapes.add_textbox(Inches(6.5), Inches(2.5), Inches(5.5), Inches(1.5))
tf = algo_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Le moteur d'itineraire resout le "
run.font.size = Pt(16)
run.font.color.rgb = DARK
run2 = p.add_run()
run2.text = "probleme du voyageur de commerce (TSP)"
run2.font.size = Pt(16)
run2.font.bold = True
run2.font.color.rgb = DARK
run3 = p.add_run()
run3.text = " pour generer la trajectoire la plus economique possible."
run3.font.size = Pt(16)
run3.font.color.rgb = DARK

moteur_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.2), Inches(5.5), Inches(1.5))
moteur_shape.fill.solid()
moteur_shape.fill.fore_color.rgb = WHITE
moteur_shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
moteur_box = slide.shapes.add_textbox(Inches(6.8), Inches(4.4), Inches(5), Inches(1.3))
tf = moteur_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Moteur d'itineraires :"
p.font.size = Pt(14)
p.font.color.rgb = ORANGE
p.font.bold = True
p2 = tf.add_paragraph()
p2.text = "Nous utilisons le couplage puissant de OpenStreetMap et OSRM (Open Source Routing Machine) pour une flexibilite totale."
p2.font.size = Pt(13)
p2.font.color.rgb = DARK

emission_box = slide.shapes.add_textbox(Inches(6.5), Inches(6.0), Inches(5.5), Inches(0.8))
tf = emission_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Reduction des emissions : "
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = DARK_GREEN
run2 = p.add_run()
run2.text = "Moins de distances parcourues = moins de CO2 rejete a Douala."
run2.font.size = Pt(14)
run2.font.color.rgb = DARK

# Map placeholder
map_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.8), Inches(5), Inches(5))
map_shape.fill.solid()
map_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
map_shape.line.color.rgb = GREEN
map_shape.line.width = Pt(2)
map_text = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(4), Inches(1))
tf = map_text.text_frame
p = tf.paragraphs[0]
p.text = "\U0001F5FA Carte Dynamique"
p.font.size = Pt(20)
p.font.color.rgb = DARK_GREEN
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Itineraires optimises en temps reel"
p2.font.size = Pt(14)
p2.font.color.rgb = GRAY
p2.alignment = PP_ALIGN.CENTER

# ========== SLIDE 5: SYSTEME DE POINTS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "SYSTEME DE POINTS & GAMIFICATION")

meca_title = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.6))
tf = meca_title.text_frame
p = tf.paragraphs[0]
p.text = "Mecanique Incitative Reelle"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK

meca_desc = slide.shapes.add_textbox(Inches(0.85), Inches(2.7), Inches(5.5), Inches(1.2))
tf = meca_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "EcoTrack transforme le civisme en gains financiers directs. Les points recoltes sont echangeables en credits de communication Orange, MTN et Camtel."
p.font.size = Pt(16)
p.font.color.rgb = DARK

formula_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(4.2), Inches(5.5), Inches(1.3))
formula_shape.fill.solid()
formula_shape.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF0)
formula_shape.line.fill.background()
formula_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(5), Inches(1.1))
tf = formula_box.text_frame
p = tf.paragraphs[0]
p.text = "Calcul dynamique des points :"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
p2 = tf.add_paragraph()
p2.text = "Points = (Note x Coeff_Impact) + Bonus_Engagement"
p2.font.size = Pt(16)
p2.font.color.rgb = ORANGE

note_box = slide.shapes.add_textbox(Inches(0.85), Inches(5.8), Inches(5.5), Inches(0.6))
tf = note_box.text_frame
p = tf.paragraphs[0]
p.text = "* Les coefficients varient selon la valeur ecologique (ex: compostage, ramassage plastique)."
p.font.size = Pt(12)
p.font.color.rgb = GRAY

# 500 FCFA box
fcfa_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.5), Inches(4.5), Inches(3.5))
fcfa_shape.fill.solid()
fcfa_shape.fill.fore_color.rgb = WHITE
fcfa_shape.line.color.rgb = ORANGE
fcfa_shape.line.width = Pt(3)
fcfa_box = slide.shapes.add_textbox(Inches(7.8), Inches(2.8), Inches(4), Inches(3))
tf = fcfa_box.text_frame
p = tf.paragraphs[0]
p.text = "500 FCFA"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Palier minimal de rachat"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = DARK_GREEN
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Convertible automatiquement et envoye par notification push."
p3.font.size = Pt(13)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER

# ========== SLIDE 6: DIFFERENCIATION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "DIFFERENCIATION : NOTRE FORCE UNIQUE")

add_card(slide, Inches(0.85), Inches(2.5), Inches(3.5), Inches(3.5),
         "Cartographie Souple",
         "Adaptee au reseau routier informel de Douala et Yaounde via OpenStreetMap, la ou Google Maps echoue souvent ou coute trop cher.",
         border_color=ORANGE)

add_card(slide, Inches(4.85), Inches(2.5), Inches(3.5), Inches(3.5),
         "Gains Reels vs Badges",
         "Nous offrons du credit de communication concret au lieu de simples badges virtuels, repondant a un besoin essentiel quotidien.",
         border_color=GREEN)

add_card(slide, Inches(8.85), Inches(2.5), Inches(3.5), Inches(3.5),
         "Zero Dependance API",
         "Solution autonome et economique n'utilisant aucune cle de facturation API externe couteuse et risquant des suspensions brutales.",
         border_color=DARK_GREEN)

# ========== SLIDE 7: BUSINESS MODEL CANVAS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "BUSINESS MODEL CANVAS")

canvas_data = [
    (Inches(0.5), Inches(2.0), Inches(2.8), Inches(2.2), "PARTENAIRES", "MTN, Orange, Camtel (Operateurs), Hysacam, Communes, OpenStreetMap.", None),
    (Inches(3.5), Inches(2.0), Inches(2.8), Inches(2.2), "ACTIVITES", "Maintenance technique, routage algorithmique, moderation, distribution credits.", None),
    (Inches(6.5), Inches(2.0), Inches(3.0), Inches(2.2), "PROPOSITION VALEUR", "Pour Citoyens : Signalement ludique, recompenses directes.\nPour Collecteurs : Reduction drastique des temps de trajet et du CO2.", ORANGE),
    (Inches(9.7), Inches(2.0), Inches(3.0), Inches(2.2), "SEGMENTS", "Urbains connectes : Citoyens actifs au Cameroun.\nOrganismes : Communes urbaines, sous-traitants de collecte.", GREEN),
    (Inches(0.5), Inches(4.5), Inches(2.8), Inches(1.5), "RESSOURCES", "Algorithme TSP, base de donnees OSM, equipe technique, plateforme web/mobile.", None),
    (Inches(3.5), Inches(4.5), Inches(2.8), Inches(1.5), "CANAUX", "Application web responsive, reseaux sociaux, partenariats B2G/B2B de proximite.", None),
    (Inches(0.5), Inches(6.3), Inches(6.0), Inches(1.0), "STRUCTURE DE COUTS", "Maintenance infrastructure cloud, budget d'achat de credits telephoniques, communication locale.", None),
    (Inches(6.7), Inches(6.3), Inches(6.0), Inches(1.0), "SOURCES DE REVENUS", "Abonnements des communes (CTD), contrats pros avec entreprises de collecte, subventions climatiques.", None),
]

for (left, top, width, height, title, desc, border) in canvas_data:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(2)
    else:
        shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.width = Pt(1)

    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN

    db = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6))
    tf = db.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = DARK

# ========== SLIDE 8: ARCHITECTURE & STACK ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "ARCHITECTURE & STACK TECHNIQUE")

arch_desc = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(11), Inches(0.8))
tf = arch_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Une architecture "
run.font.size = Pt(18)
run.font.color.rgb = DARK
run2 = p.add_run()
run2.text = "robuste et evolutive"
run2.font.size = Pt(18)
run2.font.color.rgb = GREEN
run2.font.bold = True
run3 = p.add_run()
run3.text = ", concue avec les technologies modernes les plus performantes :"
run3.font.size = Pt(18)
run3.font.color.rgb = DARK

techs = [
    ("Next.js 15+", "Framework moderne et rapide pour des chargements instantanes."),
    ("Prisma & MySQL", "Base de donnees relationnelle hautement securisee."),
    ("Leaflet & OSM", "Cartographie reactive sans frais de licence."),
    ("OSRM Engine", "Calcul d'itineraire haute performance."),
]

for i, (name, desc) in enumerate(techs):
    left = Inches(0.85 + i * 3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.2), Inches(2.8), Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    card.line.width = Pt(1)

    name_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(3.8), Inches(2.4), Inches(0.5))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    desc_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(4.5), Inches(2.4), Inches(1.2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 9: L'EQUIPE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "UNE EQUIPE COMPLETE & PASSIONNEE")

team_desc = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(5.5), Inches(1.5))
tf = team_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Concepteurs Independants de A a Z"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK
p2 = tf.add_paragraph()
p2.text = ""
p3 = tf.add_paragraph()
run = p3.add_run()
run.text = "Composee de "
run.font.size = Pt(16)
run.font.color.rgb = DARK
run2 = p3.add_run()
run2.text = "4 etudiants camerounais"
run2.font.size = Pt(16)
run2.font.color.rgb = ORANGE
run2.font.bold = True
run3 = p3.add_run()
run3.text = ", notre equipe unit ses forces pour le developpement et la mise en oeuvre de la solution a Douala."
run3.font.size = Pt(16)
run3.font.color.rgb = DARK

quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(4.5), Inches(5.5), Inches(1.2))
quote_shape.fill.solid()
quote_shape.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xE8)
quote_shape.line.fill.background()
quote_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.7), Inches(5), Inches(1))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Innover pour preserver la sante publique et notre environnement urbain camerounais."'
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = DARK

members = [
    ("K", "Klein", "Developpeur Lead", ORANGE),
    ("E", "Eunice", "UX/UI Designer", DARK_GREEN),
    ("L", "Lethicia", "Data Analyst", GREEN),
    ("P", "Parker", "Product Manager", ORANGE),
]

positions = [(Inches(7.5), Inches(1.8)), (Inches(10), Inches(1.8)), (Inches(7.5), Inches(4.5)), (Inches(10), Inches(4.5))]

for i, ((letter, name, role, color), (left, top)) in enumerate(zip(members, positions)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), top + Inches(0.2), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_GREEN
    circle.line.color.rgb = color
    circle.line.width = Pt(3)

    letter_box = slide.shapes.add_textbox(left + Inches(0.6), top + Inches(0.35), Inches(1), Inches(0.7))
    tf = letter_box.text_frame
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    name_box = slide.shapes.add_textbox(left, top + Inches(1.3), Inches(2.2), Inches(0.4))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    role_box = slide.shapes.add_textbox(left, top + Inches(1.7), Inches(2.2), Inches(0.4))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(12)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 10: VISION D'EXPANSION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_green_bar(slide)
add_title_text(slide, "VISION D'EXPANSION REGIONALE")

roadmap = [
    ("Q3 2026", "Phase pilote de lancement a Douala (arrondissement de Douala 5eme)."),
    ("Q1 2027", "Deploiement generalise et partenariats avec la ville de Yaounde."),
    ("Q4 2027", "Lancement dans les villes secondaires du Cameroun (Garoua, Bafoussam)."),
    ("Horizon 2028", "S'etendre a 5 grandes metropoles d'Afrique centrale."),
]

for i, (period, desc) in enumerate(roadmap):
    left = Inches(0.85 + i * 3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.0), Inches(2.8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF0, 0xFA, 0xF0)
    card.line.color.rgb = RGBColor(0xDD, 0xEE, 0xDD)
    card.line.width = Pt(1)

    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.15), Inches(3.1), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN if i % 2 == 1 else ORANGE
    dot.line.fill.background()

    period_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(3.6), Inches(2.4), Inches(0.5))
    tf = period_box.text_frame
    p = tf.paragraphs[0]
    p.text = period
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    desc_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 11: CONCLUSION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x1A)

add_green_bar(slide, top=Inches(1.0), height=Inches(0.5))
conc_title = slide.shapes.add_textbox(Inches(0.85), Inches(0.9), Inches(5), Inches(0.6))
tf = conc_title.text_frame
p = tf.paragraphs[0]
p.text = "CONCLUSION"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

quote_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.5), Inches(6), Inches(2.5))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Cet enfant qui ne peut pas traverser sa rue certains matins a cause des ordures -- avec EcoTrack, demain il traverse."'
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE

slogan_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.5), Inches(6), Inches(0.8))
slogan_shape.fill.solid()
slogan_shape.fill.fore_color.rgb = DARK_GREEN
slogan_shape.line.fill.background()
slogan_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.6), Inches(5.5), Inches(0.6))
tf = slogan_box.text_frame
p = tf.paragraphs[0]
p.text = "EcoTrack, vers une ville plus propre."
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Why we deserve to win
why_box = slide.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))
tf = why_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Pourquoi EcoTrack merite de gagner :"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = GREEN

items = [
    "Innovation reelle et differenciante",
    "Impact environnemental mesurable",
    "Faisabilite technique prouvee",
    "Business model viable",
    "Equipe complementaire et motivee",
]
for item in items:
    p2 = tf.add_paragraph()
    p2.text = f"  {item}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(8)

p3 = tf.add_paragraph()
p3.text = ""
p4 = tf.add_paragraph()
p4.text = '"Ensemble, transformons nos villes en modeles d\'ecologie citoyenne !"'
p4.font.size = Pt(13)
p4.font.italic = True
p4.font.color.rgb = YELLOW

output_path = "/Users/mac/devops-exercices/EcoTrack_Presentation.pptx"
prs.save(output_path)
print(f"Presentation generee : {output_path}")
